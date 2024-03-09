from pickle import INST
from turtle import clear
import pandas as pd
from pathlib import Path
import os
from re import split
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from copy import deepcopy

# Get the path to the user's desktop directory
desktop_path = Path.home() / "Desktop"

folder_name = "Scout"
folder_path = desktop_path / (folder_name)

# Check if the folder already exists
if os.path.exists(folder_path):
    l = "Exists"
else:
    # Create the folder
    os.mkdir(folder_path)



def instructions():
    print("\t\t......Scouting Program 2.0 by Klei Nagode......")
    print("                             INSTRUCTIONS FOR USE")
    print("\nA folder on your desktop will appear named Scout. Put the Powerpoint presentation")
    print("that you would want to update in that folder and follow the instructions on the screen")
    
    stop = input("\nPress Enter to move on...\n")
    os.system("cls")
    
def url_modification():
    
    url = "https://naiastats.prestosports.com/sports/mbkb/2023-24/conf/Sooner/teams/untdallas?tmpl=teaminfo-network-monospace-template&sort=ptspg"

    print("1. SAGU\n2. Science and Arts\n3. Langston\n4. Wayland\n5. Central Christian\n6. Mid-America Christian")
    print("7. Oklahoma City\n8. UNT Dallas\n9. Southwestern Christian\n10. John Brown\n11. Texas Wesleyan\n12. Oklahoma Panhandle State")
    
    cho = int(input("\nEnter the choice for the team you want stats for: "))
    
    if cho == 1:
        new_url = url.replace("untdallas","sagu")
    elif cho == 2:
        new_url = url.replace("untdallas","scienceandarts")
    elif cho == 3:
        new_url = url.replace("untdallas","langston")
    elif cho == 4:
        new_url = url.replace("untdallas","waylandbaptisttx")
    elif cho == 5:
        new_url = url.replace("untdallas","centralchristianks")
    elif cho == 6:
        new_url = url.replace("untdallas","midamericachristian")
    elif cho == 7:
        new_url = url.replace("untdallas","oklahomacity")
    elif cho == 8:
        new_url = url.replace("untdallas","untdallas")
    elif cho == 9:
        new_url = url.replace("untdallas","southwesternchristian")
    elif cho == 10:
        new_url = url.replace("untdallas","johnbrown")
    elif cho == 11:
        new_url = url.replace("untdallas","texaswesleyan")
    elif cho == 12:
        new_url = url.replace("untdallas","oklahomapanhandlestate")


    new_url = new_url.replace("2023-24",input("Enter the season you want the stats for in this format 2023-24: "))

    return new_url
    
def table_finder_modification():
   
   url = url_modification()
   
   df = pd.read_html(url)
   
   raw_table = df[1]
   conference_table = raw_table
  
   new_columns = [col[1] for col in raw_table.columns]

   # Rename the columns with the new names
   conference_table.columns = new_columns
   
   raw_table["#"] = raw_table["#"].fillna("*") #Fills the NaN values with stars
  
   count = 0
   for index, item in raw_table["#"].items(): #After it finds two stars in a row it saves the index value
       if item == "*":
           count += 1
       else:
           count = 0
           
       if count == 2:
           location = index
           break       
       
   conference_table = raw_table.drop(raw_table.index[location:]) #Here it removes the rows after the index vlaue
   
   conference_table = conference_table.drop(conference_table.index[0]) #Drops the 0 row
   
   conference_table = conference_table.replace("-", 0) #Replaces - with 0
   
   
   #Finds the conference games and all games with index and stars
   co = 0
   for ind, it in raw_table["#"].items():
       if it == "*":
           co += 1
       else:
           co = 0
           
       if co == 2:
           loca = ind
           break
   
   all_games_played = float(raw_table.loc[loca +1 , "GP"])
   con_games_played = float(raw_table.loc[loca + 2, "GP"])
      
   
   conference_table["GP"] = conference_table["GP"].astype(int) #Rounds all the values to a whole num
   conference_table["GS"] = conference_table["GS"].astype(int)
   conference_table["GP"] = conference_table["GP"].astype(str) #Converts all the needed columns to string 
   conference_table["GS"] = conference_table["GS"].astype(str)
   
   for index, item in conference_table["GP"].items():
        conference_table.loc[index,"GP"]= item + "-"        
   
   conference_table["GP-GS"] = conference_table["GP"] + conference_table["GS"] #Combines the GP and GS Columns to look like this: GP-GS
   
   del conference_table["GS"] #Deleting the GP Column
   
   conference_table["GP"] = conference_table["GP-GS"] #Replaces the GP column with the GP-GS one

   del conference_table["GP-GS"] #Deleting the duplicated GP-GS column
   
   conference_table.rename(columns = {"GP" : "GP-GS"}, inplace = True) #Renames the GP column to GP-GS3
   
   del conference_table["PTS"]
   del conference_table["STL/G"]
   del conference_table["BLK/G"]
   del conference_table["TO/G"]
   del conference_table["A/G"]
   del conference_table["A"]   #Removes the unwanted columns
   del conference_table["STL"]
   del conference_table["BLK"]
   del conference_table["TO"]
   del conference_table["DQ"]
   del conference_table["PF"]
   del conference_table["TOT"]
   del conference_table["DEF"]
   del conference_table["OFF"]
   del conference_table["MIN"]
   
   #Renaming all the columns
   conference_table.columns = ["#", "Player", "GP-GS", "MIN", "FG Att/Game", "FG%", "FG Att are 3s", "3FG%", "FT Att/Game", "FT%", "RPG", "ASST/TO", "PPG"]

   #Spliting the values with "-"
   conference_table["FG Att/Game"] = conference_table["FG Att/Game"].str.split("-").str[-1]
   conference_table["FG Att are 3s"] = conference_table["FG Att are 3s"].str.split("-").str[-1]
   conference_table["FT Att/Game"] = conference_table["FT Att/Game"].str.split("-").str[-1]
   

   drop_list1 = [] #Stores the rows that need to be deleted
   drop_list2 = []
   player_num_list = []
   player_name_list = []
   
       
   for ind, item in conference_table["#"].items():
           
       if item == "*":
           drop_list1.append(ind-1)
       
   all_table = conference_table.drop(conference_table.index[drop_list1])    
       
   for ind, item in conference_table["#"].items():

        if item != "*":
               
            drop_list2.append(ind-1)
            player_num_list.append(int(item))
            name = conference_table.loc[ind, "Player"]  #Saves the name in a variable and strips the . from it, also it adds it to the names list
            player_name_list.append(name.replace(".", ""))
               
       
   conference_table = conference_table.drop(conference_table.index[drop_list2])       
       
   all_table["Player"] = player_name_list
   conference_table["#"] = player_num_list   #Renames the empty columns with numbers and names
   conference_table["Player"] = player_name_list
   
   conference_table["FG Att are 3s"] = conference_table["FG Att are 3s"].astype(float)
   conference_table["FG Att/Game"] = conference_table["FG Att/Game"].astype(float)  
   conference_table["FT Att/Game"] =  conference_table["FT Att/Game"].astype(float)

   #Calculations for conference
   conference_table["FG Att are 3s"] = ((conference_table["FG Att are 3s"] / conference_table["FG Att/Game"]) * 100).round()
   conference_table["FG Att/Game"] = (conference_table["FG Att/Game"] / con_games_played).round(1)
   conference_table["FT Att/Game"] = (conference_table["FT Att/Game"] / con_games_played).round(1)
   
   #Checking for division by 0
   values = conference_table["FG Att are 3s"]
   values = values.replace(np.nan, 0)
   conference_table["FG Att are 3s"] = values.astype(str) + "%"
   
   conference_table["FT%"] = conference_table["FT%"].astype(str) + "%" #Adding the % sign to FT and #FG
   conference_table["3FG%"] = conference_table["3FG%"].astype(str) + "%"
   conference_table["FG%"] = conference_table["FG%"].astype(str) + "%"
   
    
   
   all_table["FG Att are 3s"] = all_table["FG Att are 3s"].astype(float)
   all_table["FG Att/Game"] = all_table["FG Att/Game"].astype(float)  
   all_table["FT Att/Game"] =  all_table["FT Att/Game"].astype(float)


   #Calculations for all games
   all_table["FG Att are 3s"] = ((all_table["FG Att are 3s"] / all_table["FG Att/Game"]) * 100).round()
   all_table["FG Att/Game"] = (all_table["FG Att/Game"] / all_games_played).round(1)
   all_table["FT Att/Game"] = (all_table["FT Att/Game"] / all_games_played).round(1)


   
   #Checking for division by 0
   values2 = all_table["FG Att are 3s"]
   values2 = values2.replace(np.nan, 0)
   all_table["FG Att are 3s"] = values2.astype(str) + "%"
   
   all_table["FT%"] = all_table["FT%"].astype(str) + "%"
   all_table["3FG%"] = all_table["3FG%"].astype(str) + "%"
   all_table["FG%"] = all_table["FG%"].astype(str) + "%"
    
   #Exports to excel
   path1 = str(folder_path / ("raw_table.xlsx"))
   raw_table.to_excel(path1)     
   path2 = str(folder_path / ("conference_stats.xlsx"))
   conference_table.to_excel(path2)
   path3 = str(folder_path / ("all_stats.xlsx"))
   all_table.to_excel(path3)  
 
def PowerPoint_creation():
    
    path2 = str(folder_path / ("conference_stats.xlsx"))
    path3 = str(folder_path / ("all_stats.xlsx")) 
    
    cho = int(input("\nWhat stats do you need to update the PowerPoint with?\n1. ALL STATS\n2. CONFERENCE STATS\n"))
   
    if cho == 1:
        df_excel = pd.read_excel(path3)
    elif cho == 2:
        df_excel = pd.read_excel(path2)
        
    presentation_name = input("Enter the Powerpoint presentation name: ")
    path4 = str(folder_path / (presentation_name + ".pptx"))
    ppt = Presentation(path4)
    
    # Iterate through slides
    for slide in ppt.slides:
    # Check if the slide contains at least two shapes and the second shape is a table
        if len(slide.shapes) > 1 and slide.shapes[1].has_table:
            table = slide.shapes[1].table  # Assuming the table is always the second shape

            # Extract the player number from the slide
            player_number = table.cell(1, 0).text

            # Find the corresponding row in the Excel DataFrame
            excel_row = df_excel[df_excel['#'] == int(player_number)]
            if not excel_row.empty:
                values = excel_row.iloc[0, 3:].tolist()
            
                # Setting the Row idx to the second row
                row_idx = 1
                row = table.rows[row_idx]
            
                # Get the font color of the cell in the first column, second row
                inherited_font_color = table.cell(1, 0).text_frame.paragraphs[0].runs[0].font.color
            
                # Update the values in the PowerPoint table
                for col_idx, cell in enumerate(row.cells):
                    if col_idx < 1:  # Skip the first three columns (indx, # and Player)
                        continue
                
                    # Clear existing content
                    new_paragraph = cell.text_frame.paragraphs[0].clear()
                
                    # Add new text with preserved formatting
                    new_run = new_paragraph.add_run()
                    new_run.text = str(values[col_idx - 1])
                
   
    
    # Save the modified PowerPoint presentation
    path5 = str(folder_path / (presentation_name + "_Updated" + ".pptx"))
    ppt.save(path5)
   

instructions()
table_finder_modification()

cho = int(input("\nDo you want to update the PowerPoint with those stats?\n1. YES\n2. NO\n"))

if cho == 1:
    PowerPoint_creation()
    
   

done = input("DONE...")
    

