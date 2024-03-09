from pickle import HIGHEST_PROTOCOL
from re import split
import numpy as np
import pandas as pd
from pathlib import Path
import os
import re

# Get the path to the user's desktop directory
desktop_path = Path.home() / "Desktop"

folder_name = "Scout"
folder_path = desktop_path / (folder_name)

#Header
print("\n\t\t-----Scout-Generation: Nagode Klei---------\n\n")

# Check if the folder already exists
if os.path.exists(folder_path):
    l = "Exists"
else:
    # Create the folder
    os.mkdir(folder_path)

#input("Enter the path to the .txt file: ")
print("1.) Put the scout Powerpoint which you want to update in the Folder: Scout on your DESKTOP")
print("2.) Go in that Folder and Create a new TEXT document in which you paste ALL the stats from the PDF (include the # row too)")
filename = input("Type in the the name of your text file: ")
path = str(folder_path / (filename + ".txt"))

Games_Played = float(input("Enter the ammount of games played: "))
print("\n")

input_txt_path = path

# Open the input file for reading
with open(input_txt_path, 'r') as input_file:
    # Read the content of the input file
    content = input_file.read()

# Remove extra spaces after commas
content = content.replace(', ', ',').replace(" Sr. ", ",Sr.").replace(" Jr.", ",Jr.").replace(" I", ",I").replace(" II", ",II").replace(" III", "III").replace(" IV", "IV")

# Open the output file for writing
with open(input_txt_path, 'w') as output_file:
    # Write the modified content to the output file
    output_file.write(content)
    



#Creates tables from the txt file
df1 = pd.read_table(path, sep=" " )
df2 = pd.read_table(path, sep=" " )


#Removes the commas in the name Column
for index, value in df1["MIN"].items(): 
    df1.loc[index, "Player"] = (df1.loc[index, "Player"].replace(",", " ")) #Removes the commas in the names
    
  
   

#Spliting the values with "-"
df2["FG-FGA"] = df2["FG-FGA"].str.split("-").str[-1]
df2["3FG-FGA"] = df2["3FG-FGA"].str.split("-").str[-1]
df2["FT-FTA"] = df2["FT-FTA"].str.split("-").str[-1]

#renaming and inserting columns into table 1
df1.rename(columns = {"TOT" : "RPG"}, inplace = True)
df1.rename(columns = {"AVG.2" : "PPG"}, inplace = True)

#Converting elements into float
df2["FG-FGA"] = df2['FG-FGA'].astype(float)
df2["FG%"] = df2['FG%'].astype(float)
df2["3FG-FGA"] = df2['3FG-FGA'].astype(float)
df2["3FG%"] = df2['3FG%'].astype(float)
df2["MIN"] = df2['MIN'].astype(float)
df2["FT-FTA"] = df2['FT-FTA'].astype(float)
df2["FT%"] = df2['FT%'].astype(float)
df2["TOT"] = df2['TOT'].astype(float)
df2["A"] = df2['A'].astype(float)
df2["TO"] = df2['TO'].astype(float)
df2["TO"] = df2['TO'].astype(float)

df1["FG%"] = df1['FG%'].astype(float)
df1["FT%"] = df1['FT%'].astype(float)
df1["RPG"]= df1['RPG'].astype(float)
df1["PPG"] = df1['PPG'].astype(float)
df1["#"] = df1['#'].astype(int)


#Deleting unwanted columns in table 1
del df1["MIN"]
del df1["FG-FGA"]
del df1["3FG-FGA"]
del df1["FT-FTA"]
del df1["A"]
del df1["TO"]
del df1["OFF"]
del df1["DEF"]
del df1["AVG.1"]
del df1["PF"]
del df1["DQ"]
del df1["BLK"]
del df1["STL"]
del df1["PTS"]

df1.rename(columns = {"AVG" : "MIN"}, inplace = True) #It's here because the upper comand kept deleting it


#Calculates and insets the wanted values and formats the table
df1.insert(loc = 5, column = "FG Att are 3s", value = (df2["3FG-FGA"]/df2["FG-FGA"]*100).round(1))    
df1.insert(loc = 4, column = "FG Att/Game", value = (df2["FG-FGA"] / Games_Played).round(1))
df1.insert(loc = 8, column = "FT Att/Game", value = (df2["FT-FTA"] / Games_Played).round(1))
df1.insert(loc = 11, column = "ASST/TO", value = (df2["A"] / df2["TO"]))
df1["RPG"] = (df2["TOT"] / Games_Played).round(1)
df1["FG%"] = (df1["FG%"] * 100).round(2).astype(str) + "%"
df1["3FG%"] = (df1["3FG%"] * 100).round(2).astype(str) + "%"
df1["FT%"] = (df1["FT%"] * 100).round(2).astype(str) + "%"


#Checking for division by 0
values = df1["FG Att are 3s"]
values = values.replace(np.nan, 0)
df1["FG Att are 3s"] = values.round(1).astype(str) + "%"

values = df1["ASST/TO"]
values = values.replace(np.nan, 0)
df1["ASST/TO"] = values.round(2).astype(str)


#Exports them into excel
path2 = str(folder_path / ("table1.xlsx"))
path3 = str(folder_path / ("table2.xlsx"))
df1.to_excel(path2)
df2.to_excel(path3)


#_________________________________________________________________________________________________________________

#PowerPoint Creation
from pptx import Presentation
from pptx.util import Inches, Pt
from copy import deepcopy


df_excel = pd.read_excel(path2)

presentation_name = input("Enter the Powerpoint presentation name: ")
path4 = str(folder_path / (presentation_name + ".pptx"))
ppt = Presentation(path4)


# Iterate through slides
for slide in ppt.slides:
    # Check if the slide contains at least two shapes and the second shape is a table
    if len(slide.shapes) > 1 and slide.shapes[1].has_table:
        table = slide.shapes[1].table  # Assuming the table is always the second shape

        # Extract the player number from the slide
        player_number = table.cell(1, 0).text

        # Find the corresponding row in the Excel DataFrame
        excel_row = df_excel[df_excel['#'] == int(player_number)]
        if not excel_row.empty:
            values = excel_row.iloc[0, 3:].tolist()
            
            # Setting the Row idx to the second row
            row_idx = 1
            row = table.rows[row_idx]
            
            # Get the font color of the cell in the first column, second row
            inherited_font_color = table.cell(1, 0).text_frame.paragraphs[0].runs[0].font.color
            
            # Update the values in the PowerPoint table
            for col_idx, cell in enumerate(row.cells):
                if col_idx < 1:  # Skip the first three columns (indx, # and Player)
                    continue
                
                # Preserve existing formatting of the cell
                existing_paragraph = cell.text_frame.paragraphs[0]
                existing_runs = existing_paragraph.runs
                existing_font = existing_runs[0].font
                
                # Clear existing content
                new_paragraph = cell.text_frame.paragraphs[0].clear()
                
                # Add new text with preserved formatting
                new_run = new_paragraph.add_run()
                new_run.text = str(values[col_idx - 1])
                
                # Preserve font properties (bold, italic, underline, size, name)
                new_run.font.bold = existing_font.bold
                new_run.font.italic = existing_font.italic
                new_run.font.underline = existing_font.underline
                new_run.font.size = existing_font.size
                new_run.font.name = existing_font.name


# Save the modified PowerPoint presentation
path5 = str(folder_path / (presentation_name + "Updated" + ".pptx"))
ppt.save(path5)

#Removing the table files
os.remove(path3)

print("\n\nUpdate Successful")
input("Press Enter to continue...")